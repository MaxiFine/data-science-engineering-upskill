from pptx import Presentation


def build_presentation():
    prs = Presentation()

    # Slide 1 – Title & Objective
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Financial Performance Review – [Period]"
    subtitle = slide.placeholders[1]
    subtitle.text = "Prepared for: CFO\nObjective: Review performance, highlight risks/opportunities, and agree 90‑day action plan."

    # Helper to add title & content slides
    def add_bullet_slide(title, bullets):
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = bullets[0]
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    # Slide 2 – Agenda
    add_bullet_slide(
        "Agenda",
        [
            "Financial overview",
            "KPIs & trends",
            "Insights & action plan",
        ],
    )

    # Slide 3 – Financial Overview: Headline Results
    add_bullet_slide(
        "Financial Overview: Headline Results",
        [
            "Revenue and gross margin overview (vs last period/target).",
            "Strongest performance: Government & Small Business segments (Paseo, Montana, Carretera, VTT).",
            "Main pressure: Enterprise deals with high discounts on VTT/Velo.",
        ],
    )

    # Slide 4 – Segment & Product Profitability
    add_bullet_slide(
        "Segment & Product Profitability",
        [
            "Chart: Profit or margin % by segment (Gov / SMB / Enterprise / Channel).",
            "Chart: Top products by profit (Paseo, Montana, Carretera, VTT).",
            "Highlight: Green = Gov & SMB; Red = Enterprise high-discount deals.",
        ],
    )

    # Slide 5 – Country Performance Overview
    add_bullet_slide(
        "Country Performance Overview",
        [
            "Chart: Margin by country × segment (focus on Canada, France).",
            "Canada: Strong Government performance (Paseo, Amarilla, Carretera).",
            "France: Strong Gov/SMB for Velo, Montana, VTT.",
        ],
    )

    # Slide 6 – KPI Dashboard
    add_bullet_slide(
        "KPI Dashboard",
        [
            "Gross margin % (overall and by segment).",
            "Share of revenue from high-margin products.",
            "Discount mix: None / Low / Medium / High.",
            "Number and value of loss-making deals.",
        ],
    )

    # Slide 7 – Trend: Discounts vs Margin
    add_bullet_slide(
        "Trend: Discounts vs Margin",
        [
            "Visual: Scatter or bar chart of discount band vs average margin.",
            "High discounts in Enterprise cluster at low or negative margins.",
            "Key message: Discount governance is the single biggest lever on profit.",
        ],
    )

    # Slide 8 – Trend: Product & Segment Mix
    add_bullet_slide(
        "Trend: Product & Segment Mix",
        [
            "Visual: Profit by product × segment over time.",
            "Profits concentrated in Gov/SMB and a few SKUs (Paseo, Montana, Carretera, VTT).",
            "Key message: We win when we sell the right SKUs into the right segments.",
        ],
    )

    # Slide 9 – Trend: COGS & Cost Drivers
    add_bullet_slide(
        "Trend: COGS & Cost Drivers",
        [
            "Visual: COGS per unit for key SKUs and variance vs prior period.",
            "Highlight low-margin SKUs with flat or rising COGS.",
            "Key message: 2–5% COGS reduction opportunity on selected SKUs.",
        ],
    )

    # Slide 10 – Insights: What Is Working
    add_bullet_slide(
        "Insights: What Is Working",
        [
            "Government & Small Business segments: stable, positive margins.",
            "Canada and France: consistently strong contributions.",
            "Products: Paseo, Montana, Carretera, VTT in Gov/SMB and Channel Partners.",
        ],
    )

    # Slide 11 – Insights: What Is Not Working
    add_bullet_slide(
        "Insights: What Is Not Working",
        [
            "Enterprise deals with High discount bands on VTT, Velo (and some Paseo/Amarilla).",
            "Loss-making or low-margin SKUs where COGS leaves no room for discounting.",
            "Inconsistent discount governance across regions/segments.",
        ],
    )

    # Slide 12 – Action Plan: Grow Profitable Sales
    add_bullet_slide(
        "Action Plan: Grow Profitable Sales",
        [
            "Implement SKU × Segment price floors and approval for High discounts.",
            "Tie discounts to volume thresholds and minimum margin.",
            "Double down on Gov/SMB in Canada and France.",
            "Prioritize Paseo, Montana, Carretera, VTT in profitable segments.",
            "Replicate Canada Channel Partner playbook for Paseo in other markets.",
        ],
    )

    # Slide 13 – Action Plan: Reduce COGS & Optimize Portfolio
    add_bullet_slide(
        "Action Plan: Reduce COGS & Optimize Portfolio",
        [
            "Consolidate suppliers for VTT/Velo components; negotiate volume discounts.",
            "Target 2–5% COGS reduction over 2 quarters for low-margin SKUs.",
            "Value-engineering and BOM simplification on loss-prone SKUs.",
            "Lean initiatives to improve yield and reduce scrap on high-volume lines.",
            "Reprice or sunset chronically negative Enterprise variants.",
        ],
    )

    # Slide 14 – 90-Day Roadmap & Risks
    add_bullet_slide(
        "90-Day Roadmap & Risks",
        [
            "Weeks 0–2: Publish price floor policy and dashboards.",
            "Weeks 3–6: Supplier negotiations and discount pilots (starting in Canada Gov).",
            "Weeks 7–12: Roll-out to France, training, and comp plan updates.",
            "Risks: Government budget cycles, supplier lock-in, enterprise churn from tougher pricing.",
            "Mitigations: Multi-year frameworks, secondary suppliers, value-based bundles and ROI proofs.",
        ],
    )

    prs.save("cfo_financial_performance_review.pptx")


if __name__ == "__main__":
    build_presentation()
